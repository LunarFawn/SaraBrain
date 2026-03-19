#!/usr/bin/env python3
"""Generate v002_user_guide.docx and v002_user_guide.pptx from content."""

import os
from pathlib import Path

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

OUT_DIR = Path(__file__).parent


# ── Word Document ──────────────────────────────────────────────────────────


def make_docx():
    doc = Document()

    # -- Styles --
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    code_style = doc.styles.add_style("CodeBlock", 1)  # paragraph style
    code_style.font.name = "Courier New"
    code_style.font.size = Pt(9)
    code_style.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    code_style.paragraph_format.space_before = Pt(4)
    code_style.paragraph_format.space_after = Pt(4)

    def heading(text, level=1):
        doc.add_heading(text, level=level)

    def para(text):
        doc.add_paragraph(text)

    def bullet(text, level=0):
        doc.add_paragraph(text, style="List Bullet")

    def code(text):
        for line in text.strip().split("\n"):
            doc.add_paragraph(line, style="CodeBlock")

    def table(headers, rows):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.style = "Light Shading Accent 1"
        for i, h in enumerate(headers):
            t.rows[0].cells[i].text = h
        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                t.rows[r_idx + 1].cells[c_idx].text = val

    # ── Title ──
    title = doc.add_heading("Sara Brain v002 — User Guide", 0)
    para("A path-of-thought brain simulation that learns, remembers, and reasons "
         "through actual neuron chains — not activation levels, not pattern matching, not forgetting.")

    # ── Philosophy ──
    heading("Philosophy & Design Reasoning")

    heading("Why path-of-thought instead of activation levels", 2)
    para("Traditional neural networks assign floating-point activation levels to nodes. "
         "When you ask 'why did you recognize apple?', the answer is an opaque matrix of weights.")
    para("Sara Brain records actual paths through neuron chains. Recognition traces real recorded paths. "
         "You can ask 'why apple' and get back the exact neuron chain that led to the conclusion, "
         "including the original teaching statement.")
    para("This also solves the false fanout problem: concept-specific relation neurons (apple_color, "
         "not fruit_color) prevent unrelated properties from leaking to all concepts in a category.")

    heading("Why no forgetting / no decay", 2)
    para("Most cognitive models implement memory decay. Sara Brain never forgets. "
         "Strength only increases via the formula 1 + ln(1 + traversals). "
         "Path similarity replaces forgetting — the overlap between downstream paths "
         "tells you how related two concepts are, without discarding anything.")

    heading("Why parallel wavefront propagation", 2)
    para("Sara Brain launches parallel wavefronts, one per input neuron, simultaneously. "
         "Each wavefront independently explores the graph. Recognition happens at intersections — "
         "neurons reached by 2+ wavefronts. This is analogous to observing commonality across "
         "multiple independent paths.")

    heading("Why concept-specific relation neurons", 2)
    para("When you teach 'an apple is red', the system creates apple_color, not fruit_color. "
         "This prevents any property reaching a shared relation from fanning out to all concepts "
         "in the same category.")

    heading("Why zero external dependencies", 2)
    para("Sara Brain's core is stdlib-only Python + SQLite. No numpy, no torch, no networkx. "
         "SQLite handles persistence, concurrency (WAL mode), and indexing. "
         "The only optional dependency is the Anthropic API for LLM translation, using urllib.request.")

    # ── Data Model ──
    heading("How the Brain Works")

    heading("Data Model", 2)

    heading("Neurons (4 types)", 3)
    table(
        ["Type", "Purpose", "Example"],
        [
            ["concept", "Subject being described", "apple, banana, circle"],
            ["property", "An attribute", "red, round, sweet"],
            ["relation", "Intermediate node linking property to concept", "apple_color, circle_shape"],
            ["association", "Groups properties under a named category", "taste, mood"],
        ],
    )

    heading("Segments (directed edges)", 3)
    para("Segments connect two neurons. Each has source_id, target_id, relation label, "
         "strength (1 + ln(1 + traversals)), traversals count, created_at, and last_used timestamps.")
    para("Strength formula: 1 + ln(1 + traversals). After 1 traversal: 1.69. After 10: 3.40. After 100: 5.62. "
         "Strength never decreases.")

    heading("Paths (provenance)", 3)
    para("A path is a recorded chain of segments representing a single learned fact. "
         "Each path stores origin_id (property), terminus_id (concept), source_text (original statement), "
         "and ordered path steps referencing segments.")

    # ── Learning Pipeline ──
    heading("Learning Pipeline", 2)
    para("When you teach a fact:")
    bullet("PARSING: Strip articles, find verb (is/are), extract subject and object, singularize")
    bullet("TAXONOMY LOOKUP: Determine property type (e.g., 'red' → 'color') and subject category")
    bullet("CHAIN BUILDING: Get or create 3 neurons (property → relation → concept) and 2 segments")
    bullet("PERSISTENCE: Record path with provenance, commit to SQLite")
    para("Key behaviors: neuron reuse (existing neurons are found, not duplicated), "
         "segment strengthening (same fact taught twice increments traversals), "
         "and concept-specific relations (every subject gets its own relation neuron).")

    # ── Recognition ──
    heading("Recognition Algorithm", 2)
    para("When you recognize:")
    bullet("Resolve input labels to neurons")
    bullet("Launch parallel wavefronts (one BFS per input neuron)")
    bullet("Find intersections — neurons reached by 2+ wavefronts")
    bullet("Strengthen all traversed segments")
    bullet("Return results sorted by confidence (number of converging paths)")

    code(
        'sara> recognize red, round\n'
        '  #1 apple (2 converging paths)\n'
        '    red → apple_color → apple\n'
        '    round → apple_shape → apple'
    )

    # ── Associations ──
    heading("Associations & Question Words")
    para("Associations let you create custom property groupings beyond the built-in taxonomy.")

    heading("Built-in associations", 2)
    table(
        ["Type", "Properties"],
        [
            ["color", "red, blue, green, yellow, orange, purple, black, white, brown, pink, crimson"],
            ["shape", "round, square, triangular, flat, oval, cylindrical, spherical"],
            ["taste", "sweet, sour, bitter, salty, savory, spicy"],
            ["texture", "smooth, rough, soft, hard, fuzzy, crunchy"],
            ["size", "big, small, large, tiny, huge"],
            ["temperature", "hot, cold, warm, cool"],
        ],
    )

    heading("Dynamic associations", 2)
    para("Create new associations at runtime with 'define' and 'describe':")
    code("sara> define mood how\nsara> describe mood as happy, sad, angry, calm")
    para("Now teaching 'a puppy is happy' creates the relation neuron puppy_mood.")

    heading("Query resolution", 2)
    para("Question word → association → properties registered under it → "
         "find paths ending at concept whose origin matches → return results.")
    code("sara> what apple color\n  apple color: red")

    # ── Categories ──
    heading("Categories")
    para("Categories are simple tags on concepts. Default is 'thing'. "
         "Built-in categories include fruit, geometric, animal, vehicle.")
    code("sara> categorize apple item\nsara> categories")

    # ── LLM ──
    heading("LLM Translation (Claude-only)")
    para("Optional natural language translation using the Anthropic Messages API. "
         "OpenAI domains are explicitly blocked. Uses urllib.request (stdlib).")
    para("Setup: 'llm set <api_key> [model]'. Usage: 'ask <question>'. "
         "The translator builds a system prompt listing available commands, sends the question to Claude, "
         "and dispatches the structured response through the REPL.")

    # ── REPL Reference ──
    heading("Complete REPL Reference")
    table(
        ["Command", "Description"],
        [
            ["teach <statement>", "Learn a fact"],
            ["recognize <inputs>", "Find concepts matching comma-separated properties"],
            ["trace <label>", "Show all outgoing paths from a neuron"],
            ["why <label>", "Show all paths leading to a neuron with provenance"],
            ["similar <label>", "Find neurons with shared downstream paths"],
            ["analyze", "Scan all neurons for path similarities"],
            ["define <name> <qword>", "Create a new association with a question word"],
            ["describe <name> as <props>", "Register properties under an association"],
            ["associations", "List all associations and their properties"],
            ["<qword> <concept> <assoc>", "Query properties via question word"],
            ["questions", "List all available question words"],
            ["categorize <concept> <cat>", "Tag a concept with a category"],
            ["categories", "List all categories"],
            ["ask <question>", "Translate natural language via Claude"],
            ["llm set <key> [model]", "Configure Claude API"],
            ["llm status / llm clear", "Check or remove LLM config"],
            ["neurons", "List all neurons"],
            ["paths", "List all recorded paths"],
            ["stats", "Show brain statistics"],
            ["save", "Force flush to disk"],
            ["quit / exit", "Exit the REPL"],
        ],
    )

    # ── Storage ──
    heading("Storage Schema")
    para("SQLite with WAL mode and foreign keys. 9 tables:")
    table(
        ["Table", "Purpose"],
        [
            ["neurons", "All neurons (concept, property, relation, association)"],
            ["segments", "Directed edges with strength and traversal counts"],
            ["paths", "Recorded fact chains with source text provenance"],
            ["path_steps", "Ordered steps within a path"],
            ["similarities", "Cached similarity analysis results"],
            ["associations", "Dynamic association → property mappings"],
            ["question_words", "Association → question word mappings"],
            ["categories", "Concept → category tags"],
            ["settings", "Key-value config (LLM settings, etc.)"],
        ],
    )
    para("Future plan: migrate to data-nut-squirrel when needs outgrow SQLite.")

    # ── Testing ──
    heading("Testing")
    para("113 tests across 13 test files. All tests use in-memory SQLite for isolation.")
    table(
        ["File", "Tests", "Coverage"],
        [
            ["test_models.py", "8", "Pure dataclass behavior"],
            ["test_storage.py", "10", "SQLite repos CRUD"],
            ["test_parser.py", "6", "Taxonomy and statement parsing"],
            ["test_learner.py", "7", "Chain creation and strengthening"],
            ["test_recognizer.py", "7", "Wavefront propagation and intersection"],
            ["test_similarity.py", "3", "Shared-path analysis"],
            ["test_integration.py", "8", "End-to-end workflows"],
            ["test_associations.py", "10", "Define, describe, list, persistence"],
            ["test_categories.py", "10", "Categories CRUD and persistence"],
            ["test_query.py", "12", "Query resolution and question words"],
            ["test_translator.py", "12", "LLM translation and domain blocking"],
        ],
    )

    out_path = OUT_DIR / "v002_user_guide.docx"
    doc.save(str(out_path))
    print(f"  Created {out_path}")


# ── PowerPoint ─────────────────────────────────────────────────────────────


def make_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_content_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, text in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = text
            else:
                p = tf.add_paragraph()
                p.text = text
            para = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.paragraphs[-1]
            para.font.size = PptxPt(18)

    def add_two_column_slide(title, left_bullets, right_bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        all_items = left_bullets + [""] + right_bullets
        for i, text in enumerate(all_items):
            if i == 0:
                tf.paragraphs[0].text = text
            else:
                p = tf.add_paragraph()
                p.text = text
            para = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.paragraphs[-1]
            para.font.size = PptxPt(16)

    # Slide 1: Title
    add_title_slide(
        "Sara Brain v002",
        "Path-of-thought brain simulation\nLearn \u2022 Remember \u2022 Reason"
    )

    # Slide 2: Core Philosophy
    add_content_slide("Core Philosophy", [
        "\u2022 Path-of-thought, not activation levels",
        "    Recognition traces real recorded paths \u2014 every conclusion is explainable",
        "",
        "\u2022 Never forgets",
        "    Strength only increases (1 + ln(1 + traversals)). Path similarity replaces decay",
        "",
        "\u2022 Parallel wavefront propagation",
        "    Multiple wavefronts launch simultaneously; intersections = recognition",
    ])

    # Slide 3: Data Model
    add_content_slide("Data Model: Neurons, Segments, Paths", [
        "Neurons (4 types):",
        "    concept (apple) \u2022 property (red) \u2022 relation (apple_color) \u2022 association (taste)",
        "",
        "Segments (directed edges):",
        "    source \u2192 target with strength = 1 + ln(1 + traversals)",
        "    Strength never decreases. Traversals only increase.",
        "",
        "Paths (provenance):",
        "    Recorded chains: property \u2192 relation \u2192 concept",
        "    Each stores the original teaching statement",
    ])

    # Slide 4: Learning Pipeline
    add_content_slide("Learning Pipeline", [
        "teach an apple is red",
        "",
        "1. PARSING: strip articles, find verb, extract subject/object",
        "2. TAXONOMY: red \u2192 color, apple \u2192 fruit \u2192 relation = apple_color",
        "3. CHAIN: create/reuse property \u2192 relation \u2192 concept neurons",
        "4. SEGMENTS: red \u2192 apple_color (has_color), apple_color \u2192 apple (describes)",
        "5. PATH: record with source_text provenance, commit to SQLite",
        "",
        "Same fact taught again \u2192 segments strengthened, never duplicated",
    ])

    # Slide 5: Recognition Algorithm
    add_content_slide("Recognition: Parallel Wavefronts", [
        "recognize red, round",
        "",
        "Wavefront 1 (red):   red \u2192 apple_color \u2192 apple",
        "Wavefront 2 (round): round \u2192 apple_shape \u2192 apple",
        "",
        "Intersection at 'apple' \u2192 RECOGNIZED (confidence: 2)",
        "",
        "Concept-specific relations prevent false fanout:",
        "  apple_color is private to apple \u2014 no leaking to banana",
        "",
        "All traversed segments strengthened after recognition",
    ])

    # Slide 6: Associations & Question Words
    add_content_slide("Associations & Question Words", [
        "Built-in: color, shape, taste, texture, size, temperature",
        "",
        "Dynamic associations:",
        "    define mood how",
        "    describe mood as happy, sad, angry, calm",
        "",
        "Query resolution:",
        "    what apple color \u2192 finds paths ending at apple with color properties",
        "",
        "Question words: what (color, shape, size), how (taste, texture, temperature)",
    ])

    # Slide 7: LLM Translation
    add_content_slide("LLM Translation (Claude-only)", [
        "Optional natural language layer using Anthropic Messages API",
        "",
        "Setup: llm set sk-ant-... claude-sonnet-4-20250514",
        "Usage: ask what does an apple taste like?",
        "       \u2192 how apple taste \u2192 apple taste: sweet",
        "",
        "Design: OpenAI domains explicitly blocked",
        "        stdlib urllib.request \u2014 no SDK dependency",
        "        Zero-temperature for deterministic translation",
        "        System prompt lists all available structured commands",
    ])

    # Slide 8: REPL Commands
    add_content_slide("REPL Command Reference", [
        "Learning:     teach <statement>",
        "Recognition:  recognize <inputs>",
        "Exploration:  trace \u2022 why \u2022 similar \u2022 analyze",
        "Associations: define \u2022 describe \u2022 associations",
        "Queries:      <qword> <concept> <assoc> \u2022 questions",
        "Categories:   categorize \u2022 categories",
        "LLM:          ask \u2022 llm set/status/clear",
        "Inspection:   neurons \u2022 paths \u2022 stats",
        "System:       save \u2022 quit/exit",
    ])

    # Slide 9: Storage & Testing
    add_content_slide("Storage & Testing", [
        "Storage: SQLite with WAL mode, foreign keys",
        "    9 tables: neurons, segments, paths, path_steps, similarities,",
        "    associations, question_words, categories, settings",
        "    Future: migrate to data-nut-squirrel",
        "",
        "Testing: 113 tests across 13 files",
        "    Models \u2022 Storage \u2022 Parser \u2022 Learner \u2022 Recognizer",
        "    Similarity \u2022 Integration \u2022 Associations \u2022 Categories",
        "    Queries \u2022 Translator",
        "",
        "Zero external dependencies (core). Python 3.11+ only.",
    ])

    out_path = OUT_DIR / "v002_user_guide.pptx"
    prs.save(str(out_path))
    print(f"  Created {out_path}")


if __name__ == "__main__":
    make_docx()
    make_pptx()
    print("  Done.")
